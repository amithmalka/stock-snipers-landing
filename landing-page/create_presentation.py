"""MGS Premium Presentation — Luxury edition with proper Hebrew RTL."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import os

CHARTS = "/Users/amitmalka/siel-app/landing-page/charts"

def create():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    GOLD = RGBColor(212, 175, 55)
    BLACK = RGBColor(5, 5, 5)
    SURFACE = RGBColor(15, 15, 15)
    SURFACE2 = RGBColor(22, 22, 22)
    WHITE = RGBColor(245, 245, 245)
    MUTED = RGBColor(140, 140, 140)
    DIM = RGBColor(65, 65, 65)
    RED = RGBColor(248, 113, 113)
    GREEN = RGBColor(52, 211, 153)
    BLUE = RGBColor(96, 165, 250)
    PURPLE = RGBColor(167, 139, 250)
    ORANGE = RGBColor(245, 158, 11)
    PINK = RGBColor(244, 114, 182)

    W = prs.slide_width
    H = prs.slide_height
    sc = [0]

    def set_rtl(paragraph):
        """Set RTL on a paragraph via XML."""
        pPr = paragraph._p.get_or_add_pPr()
        pPr.set('rtl', '1')

    def set_bg(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BLACK

    def rect(slide, l, t, w, h, color):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
        return s

    def gold_bars(slide):
        rect(slide, 0, 0, W, Inches(0.04), GOLD)
        rect(slide, 0, H - Inches(0.04), W, Inches(0.04), GOLD)

    def gold_line(slide, l, t, w):
        rect(slide, l, t, w, Inches(0.02), GOLD)

    def dim_line(slide, l, t, w):
        rect(slide, l, t, w, Inches(0.01), DIM)

    def logo(slide):
        tb = slide.shapes.add_textbox(Inches(11), Inches(0.2), Inches(2), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = "MGS"; r.font.color.rgb = GOLD; r.font.size = Pt(13); r.font.bold = True
        p.alignment = PP_ALIGN.RIGHT

    def snum(slide):
        sc[0] += 1
        tb = slide.shapes.add_textbox(Inches(0.4), H - Inches(0.35), Inches(0.6), Inches(0.2))
        p = tb.text_frame.paragraphs[0]
        p.text = str(sc[0]); p.font.color.rgb = DIM; p.font.size = Pt(9)

    def base(slide):
        set_bg(slide); gold_bars(slide); logo(slide); snum(slide)

    def txt(slide, l, t, w, h, text, sz=22, col=WHITE, bold=False, align=PP_ALIGN.RIGHT, rtl=True):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        if rtl:
            set_rtl(p)
        r = p.add_run(); r.text = text; r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
        p.alignment = align
        return tb

    def txt_multi(slide, l, t, w, h, lines, sz=18, col=WHITE, line_spacing=1.5):
        """Multiple lines with proper RTL per line."""
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            set_rtl(p)
            r = p.add_run(); r.text = line; r.font.size = Pt(sz); r.font.color.rgb = col
            p.alignment = PP_ALIGN.RIGHT
            p.space_after = Pt(sz * 0.6)
        return tb

    def card(slide, l, t, w, h, title, desc_lines, tcol=GOLD, acol=GOLD):
        # Card background
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = SURFACE2
        s.line.color.rgb = RGBColor(38, 38, 38); s.line.width = Pt(0.75)
        # Gold accent top
        rect(slide, l + Inches(0.4), t + Inches(0.15), w - Inches(0.8), Inches(0.02), acol)
        # Title
        txt(slide, l + Inches(0.25), t + Inches(0.35), w - Inches(0.5), Inches(0.45),
            title, sz=16, col=tcol, bold=True, align=PP_ALIGN.CENTER)
        # Description
        if isinstance(desc_lines, str):
            desc_lines = desc_lines.split('\n')
        txt_multi(slide, l + Inches(0.25), t + Inches(0.9), w - Inches(0.5), h - Inches(1.1),
                  desc_lines, sz=13, col=MUTED)

    def chart_img(slide, name, l, t, w, h=None):
        path = f"{CHARTS}/{name}.png"
        if os.path.exists(path):
            if h:
                slide.shapes.add_picture(path, l, t, w, h)
            else:
                slide.shapes.add_picture(path, l, t, w)

    def bullet_slide(slide, title, badge, points, point_sz=20):
        """Full content slide with title, badge, and RTL bullet points."""
        txt(slide, Inches(1), Inches(0.7), Inches(10.5), Inches(0.6), title, sz=28, col=WHITE, bold=True)
        # Badge
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.38))
        bg.fill.solid(); bg.fill.fore_color.rgb = SURFACE2
        bg.line.color.rgb = GOLD; bg.line.width = Pt(1)
        bp = bg.text_frame.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        r = bp.add_run(); r.text = badge; r.font.color.rgb = GOLD; r.font.size = Pt(12)
        gold_line(slide, Inches(1), Inches(1.35), Inches(11))

        for i, point in enumerate(points):
            y = Inches(1.7) + Inches(i * 0.75)
            # Gold diamond
            d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(11.5), y + Inches(0.08), Inches(0.16), Inches(0.16))
            d.fill.solid(); d.fill.fore_color.rgb = GOLD; d.line.fill.background()
            txt(slide, Inches(1), y, Inches(10.3), Inches(0.55), point, sz=point_sz)

    NS = new_slide = lambda: prs.slides.add_slide(prs.slide_layouts[6])

    # ══════════════════════════════════════════════════════
    # 1 — TITLE
    # ══════════════════════════════════════════════════════
    s = NS(); set_bg(s); gold_bars(s); snum(s)

    # Decorative vertical gold bar
    rect(s, Inches(1), Inches(1.5), Inches(0.035), Inches(4.2), GOLD)
    # Subtle horizontal accents
    dim_line(s, Inches(2), Inches(2.5), Inches(9))
    dim_line(s, Inches(2), Inches(4.8), Inches(9))

    txt(s, Inches(2), Inches(1.6), Inches(9.5), Inches(1), "MGS", sz=86, col=GOLD, bold=True, align=PP_ALIGN.CENTER, rtl=False)
    txt(s, Inches(2), Inches(2.7), Inches(9.5), Inches(0.7), "Money Growth System", sz=32, col=WHITE, align=PP_ALIGN.CENTER, rtl=False)
    gold_line(s, Inches(4.8), Inches(3.5), Inches(3.5))
    txt(s, Inches(2), Inches(3.8), Inches(9.5), Inches(0.6), "תוכנית הלימודים המלאה", sz=26, col=MUTED, align=PP_ALIGN.CENTER)
    txt(s, Inches(2), Inches(5), Inches(9.5), Inches(0.5), "עמית ואליאב מלכה  ·  הדרך שלכם לחופש כלכלי", sz=16, col=DIM, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════
    # 2 — TABLE OF CONTENTS
    # ══════════════════════════════════════════════════════
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.6), "תוכנית הלימודים", sz=36, col=GOLD, bold=True)
    gold_line(s, Inches(1), Inches(1.35), Inches(11))

    modules = [
        ("0", "ההקדמה", "למה בחרנו בדרך הזו", DIM),
        ("1", "היסודות", "שוק ההון · אינפלציה · ריבית דריבית", BLUE),
        ("2", "האסטרטגיה", "S&P 500 · Nasdaq · SCHD · Bitcoin", GOLD),
        ("3", "הפרקטיקה", "ברוקר · מחשבון MGS · הזרקת הון", GREEN),
        ("4", "חופש כלכלי", "חוק 4% · תזרים · פרישה מוקדמת", PURPLE),
        ("5", "אוטומציה", "מיסוי · שגר ושכח · קהילת ליווי", ORANGE),
    ]
    for i, (num, title, sub, col) in enumerate(modules):
        y = Inches(1.7) + Inches(i * 0.88)
        # Number circle
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), y + Inches(0.05), Inches(0.48), Inches(0.48))
        c.fill.solid(); c.fill.fore_color.rgb = SURFACE2
        c.line.color.rgb = col; c.line.width = Pt(1.5)
        cp = c.text_frame.paragraphs[0]; cp.text = num; cp.font.color.rgb = col; cp.font.size = Pt(15); cp.font.bold = True; cp.alignment = PP_ALIGN.CENTER
        txt(s, Inches(2), y, Inches(8.5), Inches(0.35), f"מודול {num}  ·  {title}", sz=21, col=WHITE, bold=True)
        txt(s, Inches(2), y + Inches(0.35), Inches(8.5), Inches(0.28), sub, sz=14, col=MUTED)
        if i < 5:
            dim_line(s, Inches(2), y + Inches(0.72), Inches(9.5))

    # ══════════════════════════════════════════════════════
    # MODULE HEADERS
    # ══════════════════════════════════════════════════════
    def mod_header(num, title, subtitle, col=GOLD):
        s = NS(); set_bg(s); gold_bars(s); snum(s)
        # Decorative elements
        rect(s, Inches(6.2), Inches(2.2), Inches(0.8), Inches(0.01), col)
        rect(s, Inches(6.55), Inches(1.8), Inches(0.01), Inches(0.8), col)
        txt(s, Inches(2), Inches(1.6), Inches(9.5), Inches(0.6), f"מודול {num}", sz=22, col=DIM, align=PP_ALIGN.CENTER)
        gold_line(s, Inches(5.2), Inches(2.9), Inches(2.8))
        txt(s, Inches(1.5), Inches(3.1), Inches(10), Inches(0.9), title, sz=48, col=col, bold=True, align=PP_ALIGN.CENTER)
        txt(s, Inches(2), Inches(4.2), Inches(9), Inches(0.5), subtitle, sz=22, col=MUTED, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════
    # MODULE 0 — INTRODUCTION
    # ══════════════════════════════════════════════════════
    mod_header("0", "ההקדמה", "למה בחרנו בדרך הזו")

    # 0.1
    s = NS(); base(s)
    bullet_slide(s, "נעים להכיר — החזון של משפחת מלכה", "שיעור 0.1", [
        "עמית ואליאב, הורים ל-4 — החלטנו שאנחנו לא מחכים לגיל 67",
        "הפרדנו בין ״כסף של עבודה״ לבין ״כסף של עתיד״",
        "בנינו מערכת שעובדת לבד — 15 דקות בחודש ותו לא",
        "החופש שלנו: להיות ביחד בבית, לגדל את הילדים בשקט",
        "ולהיות במקום שבו אפשר לתת לאחרים — לא רק לשרוד",
    ])

    # 0.2
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(10.5), Inches(0.6), "מפת הדרכים של הקורס", sz=28, col=WHITE, bold=True)
    bg2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.38))
    bg2.fill.solid(); bg2.fill.fore_color.rgb = SURFACE2; bg2.line.color.rgb = GOLD; bg2.line.width = Pt(1)
    bp2 = bg2.text_frame.paragraphs[0]; bp2.alignment = PP_ALIGN.CENTER
    r2 = bp2.add_run(); r2.text = "שיעור 0.2"; r2.font.color.rgb = GOLD; r2.font.size = Pt(12)
    gold_line(s, Inches(1), Inches(1.35), Inches(11))

    steps = [
        ("01", "נלמד את השפה", "מושגי יסוד\nשוק ההון מאפס"),
        ("02", "נבחר את הנכסים", "4 נכסים שבונים\nאימפריה"),
        ("03", "נקים את המערכת", "ברוקר + מחשבון\nMGS הבלעדי"),
        ("04", "נתכנן חופש כלכלי", "חוק ה-4%\nופרישה מוקדמת"),
        ("05", "שגר ושכח", "15 דקות בחודש\nזה הכול"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.5 + i * 2.5)
        card(s, x, Inches(1.7), Inches(2.2), Inches(3.2), f"{num}\n{title}", desc.split('\n'))
        if i < 4:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.25), Inches(3), Inches(0.22), Inches(0.22))
            arr.fill.solid(); arr.fill.fore_color.rgb = GOLD; arr.line.fill.background()

    txt(s, Inches(1), Inches(5.3), Inches(11), Inches(0.5),
        "בלי אנגלית  ·  בלי מתמטיקה  ·  בלי ניסיון קודם  ·  רק רצון ו-15 דקות בחודש", sz=17, col=GOLD, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════
    # MODULE 1 — FOUNDATIONS
    # ══════════════════════════════════════════════════════
    mod_header("1", "היסודות", "השפה החדשה שלכם", BLUE)

    # 1.1 Concepts
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(10.5), Inches(0.6), "בורסה ושוק ההון — מאפס", sz=28, col=WHITE, bold=True)
    bg3 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.38))
    bg3.fill.solid(); bg3.fill.fore_color.rgb = SURFACE2; bg3.line.color.rgb = GOLD; bg3.line.width = Pt(1)
    bp3 = bg3.text_frame.paragraphs[0]; bp3.alignment = PP_ALIGN.CENTER
    r3 = bp3.add_run(); r3.text = "שיעור 1.1"; r3.font.color.rgb = GOLD; r3.font.size = Pt(12)
    gold_line(s, Inches(1), Inches(1.35), Inches(11))

    concepts = [
        ("שוק ההון", "זו לא הימורים — זו שותפות.\nכשאתם קונים מניה של אפל,\nאתם שותפים שלהם.\nכשהם מרוויחים — אתם מרוויחים.", GOLD),
        ("מניה", "חלק מבעלות על חברה.\nאם יש מיליון מניות ואתם\nמחזיקים 100 — אתם בעלים\nשל 0.01% מהחברה.", BLUE),
        ("מדד (Index)", "סל של מאות חברות ביחד.\nS&P 500 = 500 החברות\nהגדולות בארה״ב.\nפיזור סיכון אוטומטי.", GREEN),
        ("ETF — תעודת סל", "קונים מניה אחת בלבד\nומקבלים חשיפה למאות חברות.\nבלי לבחור מניות בעצמכם.\nהכלי המרכזי של MGS.", PURPLE),
    ]
    for i, (title, desc, col) in enumerate(concepts):
        x = Inches(0.4 + i * 3.15)
        card(s, x, Inches(1.6), Inches(2.95), Inches(4.2), title, desc.split('\n'), tcol=col, acol=col)

    txt(s, Inches(1), Inches(6.1), Inches(11), Inches(0.4),
        "עם ETF לא צריך לבחור מניות. קונים את כל השוק בלחיצת כפתור.", sz=16, col=GOLD, align=PP_ALIGN.CENTER)

    # 1.2 Inflation
    s = NS(); base(s)
    bullet_slide(s, "למה MGS — ולמה דווקא עכשיו?", "שיעור 1.2", [])

    chart_img(s, "inflation_erosion", Inches(0.4), Inches(1.6), Inches(8), Inches(4.3))

    card(s, Inches(8.8), Inches(1.7), Inches(4), Inches(1.8),
         "הכסף שלכם נשחק", ["₪100,000 בבנק היום", "= ₪74,400 בכוח קנייה", "תוך 10 שנים בלבד"],
         tcol=RED, acol=RED)
    card(s, Inches(8.8), Inches(3.7), Inches(4), Inches(2.3),
         "כוחו של זמן", ["להתחיל היום ב-₪1,000", "עדיף על להתחיל בעוד שנה ב-₪5,000", "", "ריבית דריבית לא מחכה לאף אחד."],
         tcol=GREEN, acol=GREEN)

    # 1.2b Compound interest proof
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "ההוכחה: הכוח של ריבית דריבית", sz=28, col=GOLD, bold=True)
    gold_line(s, Inches(1), Inches(1.3), Inches(11))
    chart_img(s, "compound_interest", Inches(0.3), Inches(1.5), Inches(8.5), Inches(5))
    # Hebrew captions next to chart
    card(s, Inches(9.2), Inches(1.6), Inches(3.7), Inches(1.6),
         "₪3,000 בחודש", ["הפקדה קבועה × 30 שנה", "בתשואה ממוצעת של 10.5%"],
         tcol=GOLD, acol=GOLD)
    card(s, Inches(9.2), Inches(3.4), Inches(3.7), Inches(1.5),
         "סה״כ הפקדות: ₪1.08M", ["רק את הכסף שהכנסתם.", "השאר — ריבית דריבית."],
         tcol=BLUE, acol=BLUE)
    card(s, Inches(9.2), Inches(5.1), Inches(3.7), Inches(1.5),
         "שווי תיק: ₪7.3M+", ["הרווח נוצר מעצמו.", "הזמן עושה את העבודה."],
         tcol=GREEN, acol=GREEN)

    # 1.3 Real estate vs MGS
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "MGS מול נדל״ן — השוואה אמיתית", sz=28, col=WHITE, bold=True)
    bg4 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.38))
    bg4.fill.solid(); bg4.fill.fore_color.rgb = SURFACE2; bg4.line.color.rgb = GOLD; bg4.line.width = Pt(1)
    bp4 = bg4.text_frame.paragraphs[0]; bp4.alignment = PP_ALIGN.CENTER
    r4 = bp4.add_run(); r4.text = "שיעור 1.3"; r4.font.color.rgb = GOLD; r4.font.size = Pt(12)
    gold_line(s, Inches(1), Inches(1.35), Inches(11))

    chart_img(s, "realestate_vs_mgs", Inches(0.3), Inches(1.5), Inches(8), Inches(4.3))

    card(s, Inches(8.6), Inches(1.6), Inches(4.2), Inches(1.3),
         "נזילות", ["נדל״ן: חודשים למכירה", "MGS: גישה מיידית תוך דקות"], acol=GOLD)
    card(s, Inches(8.6), Inches(3.05), Inches(4.2), Inches(1.3),
         "הון התחלתי", ["נדל״ן: מקדמה של ₪300,000+", "MGS: אפשר להתחיל מ-₪1,000"], acol=GOLD)
    card(s, Inches(8.6), Inches(4.5), Inches(4.2), Inches(1.3),
         "ניהול שוטף", ["נדל״ן: שעות של כאבי ראש", "MGS: 15 דקות בחודש — נקודה"], acol=GOLD)

    txt(s, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
        "התיק שלכם הוא ״הנדל״ן הדיגיטלי״ — בלי שוכרים, בלי תיקונים, בלי משכנתא",
        sz=15, col=GOLD, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════
    # MODULE 2 — STRATEGY
    # ══════════════════════════════════════════════════════
    mod_header("2", "האסטרטגיה", "הנכסים שמרכיבים את האימפריה")

    # Allocation overview
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "הקצאת הנכסים של MGS", sz=30, col=WHITE, bold=True)
    gold_line(s, Inches(1), Inches(1.3), Inches(11))

    chart_img(s, "allocation_pie", Inches(0.3), Inches(1.1), Inches(5.3), Inches(5.3))

    assets = [
        ("S&P 500  —  30%", ["500 החברות הכי גדולות בעולם", "VOO / SPY / IVV"], GOLD),
        ("Nasdaq  —  30%", ["טכנולוגיה, חדשנות, עתיד", "QQQ / QQQM"], PURPLE),
        ("SCHD  —  25%", ["דיבידנדים — ההכנסה הפסיבית שלכם", "SCHD / VYM / DGRO"], PINK),
        ("ביטקוין  —  15%", ["נכס הרזרבה הדיגיטלי", "IBIT / FBTC"], ORANGE),
    ]
    for i, (title, desc, col) in enumerate(assets):
        y = Inches(1.4) + Inches(i * 1.4)
        card(s, Inches(6), y, Inches(6.5), Inches(1.2), title, desc, tcol=col, acol=col)

    # 2.1 S&P 500 history + crisis proof
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "מנוע הצמיחה — 30 שנות הוכחה", sz=28, col=WHITE, bold=True)
    bg5 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.38))
    bg5.fill.solid(); bg5.fill.fore_color.rgb = SURFACE2; bg5.line.color.rgb = GOLD; bg5.line.width = Pt(1)
    bp5 = bg5.text_frame.paragraphs[0]; bp5.alignment = PP_ALIGN.CENTER
    r5 = bp5.add_run(); r5.text = "שיעור 2.1"; r5.font.color.rgb = GOLD; r5.font.size = Pt(12)
    gold_line(s, Inches(1), Inches(1.35), Inches(11))

    chart_img(s, "sp500_history", Inches(0.3), Inches(1.5), Inches(9), Inches(4.8))

    card(s, Inches(9.6), Inches(1.6), Inches(3.3), Inches(1.4),
         "תשואה ממוצעת", ["~10.5% בשנה", "לאורך 30 שנה"], acol=GOLD)
    card(s, Inches(9.6), Inches(3.2), Inches(3.3), Inches(1.5),
         "החברות שלכם", ["Apple · Microsoft", "Amazon · Google", "Nvidia · Meta"], acol=PURPLE)
    card(s, Inches(9.6), Inches(4.9), Inches(3.3), Inches(1.3),
         "60% מהתיק", ["S&P 500 + Nasdaq", "= המנועים שלכם"], acol=GREEN)

    # Crisis recovery
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "ההוכחה: השוק תמיד חוזר — תמיד", sz=28, col=GREEN, bold=True)
    gold_line(s, Inches(1), Inches(1.3), Inches(11))
    chart_img(s, "crisis_recovery", Inches(0.3), Inches(1.5), Inches(8.5), Inches(4.8))

    card(s, Inches(9.2), Inches(1.6), Inches(3.7), Inches(1.3),
         "משבר 2008", ["ירד 57% — התאושש תוך 5.5 שנים.", "מי שקנה בתחתית — הכפיל."],
         tcol=RED, acol=RED)
    card(s, Inches(9.2), Inches(3.1), Inches(3.7), Inches(1.3),
         "קורונה 2020", ["ירד 34% — חזר תוך 5 חודשים!", "ההתאוששות הכי מהירה בהיסטוריה."],
         tcol=ORANGE, acol=ORANGE)
    card(s, Inches(9.2), Inches(4.6), Inches(3.7), Inches(1.3),
         "המסקנה", ["כל ירידה הייתה הזדמנות קנייה.", "מי שנשאר — תמיד הרוויח."],
         tcol=GREEN, acol=GREEN)

    txt(s, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
        "בהיסטוריה של שוק ההון, אף משקיע לטווח ארוך לא הפסיד כסף בתקופה של 15+ שנים.",
        sz=15, col=GOLD, align=PP_ALIGN.CENTER)

    # 2.2 SCHD
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "״השכירות החודשית״ שלכם — דיבידנדים", sz=28, col=WHITE, bold=True)
    bg6 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.38))
    bg6.fill.solid(); bg6.fill.fore_color.rgb = SURFACE2; bg6.line.color.rgb = GOLD; bg6.line.width = Pt(1)
    bp6 = bg6.text_frame.paragraphs[0]; bp6.alignment = PP_ALIGN.CENTER
    r6 = bp6.add_run(); r6.text = "שיעור 2.2"; r6.font.color.rgb = GOLD; r6.font.size = Pt(12)
    gold_line(s, Inches(1), Inches(1.35), Inches(11))

    chart_img(s, "dividend_growth", Inches(0.3), Inches(1.5), Inches(8.3), Inches(4.5))

    card(s, Inches(9), Inches(1.6), Inches(3.8), Inches(1.7),
         "מה זה SCHD?", ["תעודת סל שמחלקת כסף", "כל רבעון — לחשבון שלכם.", "תשואה: ~3.5% בשנה", "צמיחת דיבידנד: ~10% בשנה"],
         acol=PINK)
    card(s, Inches(9), Inches(3.5), Inches(3.8), Inches(1.5),
         "למה זה מדהים?", ["מקבלים מזומן — בלי למכור כלום.", "ההכנסה גדלה כל שנה.", "כמו שכירות, רק בלי שוכרים."],
         tcol=GREEN, acol=GREEN)
    card(s, Inches(9), Inches(5.2), Inches(3.8), Inches(1.2),
         "25% מהתיק", ["מנוע ההכנסה הפסיבית שלכם."],
         acol=GOLD)

    # 2.3 Bitcoin
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "נכס הרזרבה — ביטקוין", sz=28, col=WHITE, bold=True)
    bg7 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.38))
    bg7.fill.solid(); bg7.fill.fore_color.rgb = SURFACE2; bg7.line.color.rgb = GOLD; bg7.line.width = Pt(1)
    bp7 = bg7.text_frame.paragraphs[0]; bp7.alignment = PP_ALIGN.CENTER
    r7 = bp7.add_run(); r7.text = "שיעור 2.3"; r7.font.color.rgb = GOLD; r7.font.size = Pt(12)
    gold_line(s, Inches(1), Inches(1.35), Inches(11))

    chart_img(s, "bitcoin_history", Inches(0.3), Inches(1.5), Inches(8.3), Inches(4.5))

    card(s, Inches(9), Inches(1.6), Inches(3.8), Inches(1.5),
         "מ-$300 ל-$95,000", ["ב-10 שנים.", "הנכס שהפך לרזרבה עולמית."],
         tcol=ORANGE, acol=ORANGE)
    card(s, Inches(9), Inches(3.3), Inches(3.8), Inches(1.4),
         "למה רק 15%?", ["מספיק ליהנות מצמיחה.", "לא יותר מדי — כדי לא לסכן."],
         acol=ORANGE)
    card(s, Inches(9), Inches(4.9), Inches(3.8), Inches(1.4),
         "גיוון חכם", ["לא תלוי בממשלות.", "הגנה בתרחישים קיצוניים."],
         acol=GREEN)

    # ══════════════════════════════════════════════════════
    # MODULE 3 — PRACTICE
    # ══════════════════════════════════════════════════════
    mod_header("3", "הפרקטיקה", "מקימים את המערכת", GREEN)

    # 3.1 Broker
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "הבית של הכסף שלכם — Colmex Pro", sz=28, col=WHITE, bold=True)
    bg8 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.38))
    bg8.fill.solid(); bg8.fill.fore_color.rgb = SURFACE2; bg8.line.color.rgb = GOLD; bg8.line.width = Pt(1)
    bp8 = bg8.text_frame.paragraphs[0]; bp8.alignment = PP_ALIGN.CENTER
    r8 = bp8.add_run(); r8.text = "שיעור 3.1"; r8.font.color.rgb = GOLD; r8.font.size = Pt(12)
    gold_line(s, Inches(1), Inches(1.35), Inches(11))

    broker_steps = [
        ("01  פתיחת חשבון", ["מדריך טכני מלא", "שלב אחרי שלב", "עם צילומי מסך"]),
        ("02  עמלות VIP", ["עמלות מופחתות", "בלעדי לתלמידי MGS", "חיסכון אדיר לאורך שנים"]),
        ("03  הפקדה ראשונה", ["העברה בנקאית פשוטה", "מאובטחת לחלוטין", "תוך דקות ספורות"]),
        ("04  הקנייה הראשונה", ["רכישת ה-ETF הראשון", "עם ליווי אישי שלנו", "צעד אחרי צעד"]),
    ]
    for i, (title, desc) in enumerate(broker_steps):
        x = Inches(0.4 + i * 3.15)
        card(s, x, Inches(1.6), Inches(2.95), Inches(3.3), title, desc, acol=GREEN)

    txt(s, Inches(1), Inches(5.3), Inches(11), Inches(0.4),
        "לא לבד — אנחנו מלווים אתכם בכל שלב. נציג אישי זמין לכל שאלה.",
        sz=16, col=GREEN, align=PP_ALIGN.CENTER)

    # Fee comparison
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "למה להשקיע לבד? ההוכחה במספרים", sz=28, col=WHITE, bold=True)
    gold_line(s, Inches(1), Inches(1.3), Inches(11))
    chart_img(s, "fee_comparison", Inches(0.3), Inches(1.5), Inches(8.5), Inches(4.8))

    card(s, Inches(9.2), Inches(1.6), Inches(3.7), Inches(1.5),
         "דמי ניהול 2%", ["נשמע מעט?", "לאורך 30 שנה זה מאות", "אלפי שקלים מהכיס שלכם."],
         tcol=RED, acol=RED)
    card(s, Inches(9.2), Inches(3.3), Inches(3.7), Inches(1.5),
         "עם MGS", ["משקיעים לבד.", "אותה תשואה — בלי לשלם", "לאף אחד על הדרך."],
         tcol=GREEN, acol=GREEN)
    card(s, Inches(9.2), Inches(5.0), Inches(3.7), Inches(1.3),
         "ההפרש", ["עשרות אלפים אחרי 20 שנה.", "מאות אלפים אחרי 30."],
         tcol=GOLD, acol=GOLD)

    txt(s, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
        "הכסף שנחסך בדמי ניהול — נשאר אצלכם וממשיך לצמוח.",
        sz=15, col=GREEN, align=PP_ALIGN.CENTER)

    # 3.2 MGS Calculator
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "מערכת MGS — המחשבון הבלעדי", sz=28, col=WHITE, bold=True)
    bg9 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.38))
    bg9.fill.solid(); bg9.fill.fore_color.rgb = SURFACE2; bg9.line.color.rgb = GOLD; bg9.line.width = Pt(1)
    bp9 = bg9.text_frame.paragraphs[0]; bp9.alignment = PP_ALIGN.CENTER
    r9 = bp9.add_run(); r9.text = "שיעור 3.2"; r9.font.color.rgb = GOLD; r9.font.size = Pt(12)
    gold_line(s, Inches(1), Inches(1.35), Inches(11))

    calc_features = [
        ("לוח בקרה חי", ["הזינו את הנתונים שלכם", "ותראו בדיוק איפה אתם", "עומדים — בלחיצה אחת."], GOLD),
        ("הזרקת הון", ["שיטה חכמה לאיזון התיק", "בלי לשלם מס מיותר.", "פשוט ויעיל."], GREEN),
        ("סימולציה עתידית", ["מריצים תחזית", "לפי תרחישים שונים.", "רואים את העתיד שלכם."], BLUE),
        ("מעקב דיבידנדים", ["רואים כמה הכנסה פסיבית", "נכנסת כל חודש.", "ההכנסה שרק גדלה."], PURPLE),
    ]
    for i, (title, desc, col) in enumerate(calc_features):
        x = Inches(0.4 + i * 3.15)
        card(s, x, Inches(1.6), Inches(2.95), Inches(3.5), title, desc, tcol=col, acol=col)

    txt(s, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
        "מחשבון שפיתחנו בעצמנו — לא קיים בשום מקום אחר. בלעדי לתלמידי MGS.",
        sz=16, col=GOLD, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════
    # MODULE 4 — FINANCIAL FREEDOM
    # ══════════════════════════════════════════════════════
    mod_header("4", "הדרך לחופש כלכלי", "מתי ואיך מגיעים לשם", PURPLE)

    # 4.1 4% Rule
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "חוק ה-4% — מתי אפשר להפסיק לעבוד?", sz=28, col=WHITE, bold=True)
    bg10 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.38))
    bg10.fill.solid(); bg10.fill.fore_color.rgb = SURFACE2; bg10.line.color.rgb = GOLD; bg10.line.width = Pt(1)
    bp10 = bg10.text_frame.paragraphs[0]; bp10.alignment = PP_ALIGN.CENTER
    r10 = bp10.add_run(); r10.text = "שיעור 4.1"; r10.font.color.rgb = GOLD; r10.font.size = Pt(12)
    gold_line(s, Inches(1), Inches(1.35), Inches(11))

    chart_img(s, "four_percent_rule", Inches(0.3), Inches(1.5), Inches(8.3), Inches(4.5))

    card(s, Inches(9), Inches(1.6), Inches(3.8), Inches(1.9),
         "מה זה חוק ה-4%?", ["מושכים 4% בשנה מהתיק.", "הכסף לא ייגמר לעולם.", "", "זה הכלל שאומר לכם:", "״הגעתם. אתם חופשיים.״"],
         tcol=PURPLE, acol=PURPLE)
    card(s, Inches(9), Inches(3.7), Inches(3.8), Inches(2.3),
         "דוגמה פשוטה", ["תיק של ₪3,000,000", "= ₪10,000 כל חודש", "", "הכנסה פסיבית.", "בלי לגעת בקרן.", "הכסף ממשיך לצמוח."],
         tcol=GREEN, acol=GREEN)

    # 4.2 Cash Flow
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "לחיות מהתזרים — הכסף עובד בשבילכם", sz=28, col=WHITE, bold=True)
    bg11 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.38))
    bg11.fill.solid(); bg11.fill.fore_color.rgb = SURFACE2; bg11.line.color.rgb = GOLD; bg11.line.width = Pt(1)
    bp11 = bg11.text_frame.paragraphs[0]; bp11.alignment = PP_ALIGN.CENTER
    r11 = bp11.add_run(); r11.text = "שיעור 4.2"; r11.font.color.rgb = GOLD; r11.font.size = Pt(12)
    gold_line(s, Inches(1), Inches(1.35), Inches(11))

    phases = [
        ("שלב 1\nצמיחה", ["שנים 1–10", "", "להפקיד בעקביות.", "לתת לריבית דריבית", "לעשות את העבודה."], BLUE),
        ("שלב 2\nמעבר", ["שנים 10–15", "", "להתחיל להעביר", "חלק מהנכסים לדיבידנדים.", "הכנסה פסיבית ראשונה."], PURPLE),
        ("שלב 3\nחופש", ["שנה 15+", "", "הרווחים מהתיק", "מכסים את כל ההוצאות.", "חופש כלכלי מלא."], GREEN),
    ]
    for i, (title, desc, col) in enumerate(phases):
        x = Inches(0.8 + i * 3.9)
        card(s, x, Inches(1.6), Inches(3.5), Inches(4), title, desc, tcol=col, acol=col)
        if i < 2:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(3.6), Inches(3.3), Inches(0.25), Inches(0.28))
            arr.fill.solid(); arr.fill.fore_color.rgb = GOLD; arr.line.fill.background()

    txt(s, Inches(1), Inches(6), Inches(11), Inches(0.4),
        "המטרה: שהכסף יעבוד בשבילכם — ולא אתם בשבילו.",
        sz=17, col=GOLD, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════
    # MODULE 5 — AUTOMATION
    # ══════════════════════════════════════════════════════
    mod_header("5", "אוטומציה וסיכום", "שגר ושכח", ORANGE)

    # 5.1 Taxes
    s = NS(); base(s)
    bullet_slide(s, "מיסוי ובירוקרטיה — בעברית פשוטה", "שיעור 5.1", [
        "מס רווח הון: 25% על הרווח בלבד — לא על מה שהפקדתם",
        "משלמים רק כשמוכרים. כל עוד לא מכרתם — אין מס. פשוט.",
        "הברוקר מנפיק דוח שנתי — מעבירים לרואה חשבון וזהו",
        "הכול מוסבר בעברית פשוטה, בלי ז׳רגון מקצועי",
    ])

    # 5.2 Set & forget
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "שגרת ״שגר ושכח״", sz=28, col=WHITE, bold=True)
    bg12 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.38))
    bg12.fill.solid(); bg12.fill.fore_color.rgb = SURFACE2; bg12.line.color.rgb = GOLD; bg12.line.width = Pt(1)
    bp12 = bg12.text_frame.paragraphs[0]; bp12.alignment = PP_ALIGN.CENTER
    r12 = bp12.add_run(); r12.text = "שיעור 5.2"; r12.font.color.rgb = GOLD; r12.font.size = Pt(12)
    gold_line(s, Inches(1), Inches(1.35), Inches(11))

    routine = [
        ("פעם בחודש\n15 דקות", ["נכנסים למערכת MGS.", "בודקים את ההקצאה.", "מפקידים לפי ההמלצה.", "סיימתם."], GOLD),
        ("פעם ברבעון\n5 דקות", ["בודקים שדיבידנדים נכנסו.", "מסתכלים על הצמיחה.", "מחייכים."], GREEN),
        ("פעם בשנה\n30 דקות", ["דוח שנתי לרואה חשבון.", "בדיקת יעדים.", "עדכון קל אם צריך."], BLUE),
    ]
    for i, (title, desc, col) in enumerate(routine):
        x = Inches(0.8 + i * 3.9)
        card(s, x, Inches(1.6), Inches(3.5), Inches(3.5), title, desc, tcol=col, acol=col)

    txt(s, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
        "לא צריך לבדוק כל יום. לא צריך להילחץ. המערכת עובדת — אתם חיים.",
        sz=17, col=ORANGE, align=PP_ALIGN.CENTER)

    # 5.3 Summary
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "סיכום — יוצאים לדרך", sz=28, col=WHITE, bold=True)
    bg13 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.38))
    bg13.fill.solid(); bg13.fill.fore_color.rgb = SURFACE2; bg13.line.color.rgb = GOLD; bg13.line.width = Pt(1)
    bp13 = bg13.text_frame.paragraphs[0]; bp13.alignment = PP_ALIGN.CENTER
    r13 = bp13.add_run(); r13.text = "שיעור 5.3"; r13.font.color.rgb = GOLD; r13.font.size = Pt(12)
    gold_line(s, Inches(1), Inches(1.35), Inches(11))

    summary_items = [
        ("למדתם את השפה", "מניה · מדד · ETF — עכשיו אתם מבינים", BLUE),
        ("בחרתם את הנכסים", "S&P 500 · Nasdaq · SCHD · ביטקוין", GOLD),
        ("הקמתם את המערכת", "ברוקר VIP + מחשבון MGS הבלעדי", GREEN),
        ("תכננתם את החופש", "חוק ה-4% + מעבר הדרגתי לדיבידנדים", PURPLE),
        ("שגר ושכח", "15 דקות בחודש — ותו לא. הכסף עובד לבד.", ORANGE),
    ]
    for i, (title, desc, col) in enumerate(summary_items):
        y = Inches(1.65) + Inches(i * 0.92)
        # Checkmark circle
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), y + Inches(0.08), Inches(0.38), Inches(0.38))
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()
        cp = c.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = "✓"; cr.font.color.rgb = BLACK; cr.font.size = Pt(14); cr.font.bold = True
        txt(s, Inches(2), y, Inches(8.5), Inches(0.35), title, sz=21, col=WHITE, bold=True)
        txt(s, Inches(2), y + Inches(0.34), Inches(8.5), Inches(0.28), desc, sz=14, col=MUTED)
        if i < 4:
            dim_line(s, Inches(2), y + Inches(0.72), Inches(9))

    txt(s, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
        "ברכת הדרך מאיתנו — אתם מוכנים. הצטרפו לקהילת MGS והתחילו את המסע.",
        sz=17, col=GOLD, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════
    # WHAT YOU GET
    # ══════════════════════════════════════════════════════
    s = NS(); base(s)
    txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.5), "מה מקבלים בקורס MGS?", sz=34, col=GOLD, bold=True)
    gold_line(s, Inches(1), Inches(1.3), Inches(11))

    benefits = [
        ("15 שיעורי וידאו", ["הדרכה צעד אחרי צעד", "מאפס ועד שגר ושכח"], GOLD),
        ("מערכת MGS בלעדית", ["מחשבון שפיתחנו בעצמנו", "לא קיים בשום מקום אחר"], BLUE),
        ("נציג אישי", ["ליווי צמוד בכל שאלה", "תמיד יש למי לפנות"], GREEN),
        ("קהילת MGS סגורה", ["קהילה של תלמידים בלבד", "שיתוף ידע ותמיכה"], PURPLE),
        ("עמלות VIP בברוקר", ["חיסכון של מאות אלפי ₪", "לאורך שנות ההשקעה"], ORANGE),
        ("גישה לכל החיים", ["כולל כל העדכונים", "ותכנים חדשים שייצאו"], PINK),
    ]
    for i, (title, desc, col) in enumerate(benefits):
        row = i // 3
        ci = i % 3
        x = Inches(0.6 + ci * 4.1)
        y = Inches(1.6) + Inches(row * 2.6)
        card(s, x, y, Inches(3.7), Inches(2.2), title, desc, tcol=col, acol=col)

    # ══════════════════════════════════════════════════════
    # FINAL SLIDE
    # ══════════════════════════════════════════════════════
    s = NS(); set_bg(s); gold_bars(s); snum(s)

    rect(s, Inches(1), Inches(1.2), Inches(0.035), Inches(4.5), GOLD)
    dim_line(s, Inches(2), Inches(2.5), Inches(9.5))
    dim_line(s, Inches(2), Inches(5.0), Inches(9.5))

    txt(s, Inches(2), Inches(1.3), Inches(9.5), Inches(1), "MGS", sz=78, col=GOLD, bold=True, align=PP_ALIGN.CENTER, rtl=False)
    txt(s, Inches(2), Inches(2.6), Inches(9.5), Inches(0.6), "Money Growth System", sz=30, col=WHITE, align=PP_ALIGN.CENTER, rtl=False)
    gold_line(s, Inches(5), Inches(3.3), Inches(3.3))
    txt(s, Inches(2), Inches(3.6), Inches(9.5), Inches(0.6), "המסע לחופש כלכלי מתחיל כאן", sz=26, col=MUTED, align=PP_ALIGN.CENTER)
    txt(s, Inches(2), Inches(4.4), Inches(9.5), Inches(0.5), "עמית ואליאב מלכה", sz=20, col=WHITE, align=PP_ALIGN.CENTER)
    txt(s, Inches(2), Inches(5.8), Inches(9.5), Inches(0.4),
        "אין באמור ייעוץ השקעות. כלי חינוכי בלבד. השקעה בשוק ההון כרוכה בסיכון.",
        sz=10, col=DIM, align=PP_ALIGN.CENTER)

    # ── Save ──
    out = "/Users/amitmalka/siel-app/landing-page/MGS_Premium_Presentation.pptx"
    prs.save(out)
    print(f"\n✅ נשמר: {out}")
    print(f"📊 סה״כ שקופיות: {sc[0]}")

if __name__ == "__main__":
    create()
